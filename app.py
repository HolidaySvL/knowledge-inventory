from pathlib import Path
from io import BytesIO
from typing import Dict, Any, List, Tuple
import zipfile
import tempfile
import subprocess
import shutil
import os
import json
import re

import pandas as pd
import streamlit as st
import requests
import plotly.express as px
from dotenv import load_dotenv

from docx import Document
from pptx import Presentation
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.utils import get_column_letter

load_dotenv()

# ============================================================
# Config
# ============================================================

SUPPORTED_EXTENSIONS = {
    ".docx", ".doc",
    ".xlsx", ".xls", ".xlsm",
    ".pptx", ".ppt",
    ".pdf",
    ".txt", ".md", ".csv",
    ".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff",
}
UPLOAD_EXTENSIONS = SUPPORTED_EXTENSIONS | {".zip"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff"}

REQUIRED_FIELDS = [
    "document_name",
    "source_path",
    "file_type",
    "document_type",
    "department",
    "business_process",
    "knowledge_category",
    "summary",
    "risk_level",
    "risk_reason",
    "recommended_action",
    "can_enter_kb",
    "needs_human_review"
]

CLASSIFY_PROMPT = """
You are an enterprise AI knowledge inventory assistant.

Your task is to analyze one business document and produce a structured inventory record.

Important rules:
- Do not rewrite the document.
- Do not invent facts.
- If something is unclear, use "Unknown" or "Needs Review".
- The output must be valid JSON only.
- Risk level must be one of: Low, Medium, High.
- can_enter_kb must be one of: Yes, No, Needs Review.
- needs_human_review must be true or false.

Metadata:
{{metadata_json}}

Document text:
{{document_text}}

Return this JSON structure exactly:
{
  "document_name": "",
  "source_path": "",
  "file_type": "",
  "document_type": "",
  "department": "",
  "business_process": "",
  "knowledge_category": "",
  "summary": "",
  "risk_level": "Low / Medium / High",
  "risk_reason": "",
  "recommended_action": "",
  "can_enter_kb": "Yes / No / Needs Review",
  "needs_human_review": true
}
"""

PROCESS_MAP_PROMPT = """
You are an enterprise knowledge management consultant.

Based on the following knowledge inventory JSON, generate a concise process map in Markdown.

Focus on:
- departments
- business processes
- related documents
- documentation gaps
- which processes are ready for knowledge-base preparation
- which processes need cleanup or human review

Inventory JSON:
{{inventory_json}}
"""

RISK_REPORT_PROMPT = """
You are an enterprise document governance analyst.

Based on the following knowledge inventory JSON, generate a concise document risk report in Markdown.

Include:
1. Executive summary
2. High-risk documents
3. Medium-risk documents
4. Documents needing human review
5. Recommended cleanup priorities
6. Priority candidates for knowledge-base preparation

Inventory JSON:
{{inventory_json}}
"""

# ============================================================
# File reading
# ============================================================

def safe_decode(file_bytes: bytes) -> str:
    for enc in ["utf-8", "utf-8-sig", "latin-1"]:
        try:
            return file_bytes.decode(enc)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("latin-1", errors="ignore")

def ocr_image_bytes(file_bytes: bytes) -> str:
    try:
        from PIL import Image
        import pytesseract
        img = Image.open(BytesIO(file_bytes))
        text = pytesseract.image_to_string(img)
        return text.strip()
    except Exception as exc:
        return f"[Image detected, but OCR is unavailable or failed: {exc}]"

def extract_docx_bytes(file_bytes: bytes, ocr_images: bool = False) -> str:
    doc = Document(BytesIO(file_bytes))
    parts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    for table_idx, table in enumerate(doc.tables, start=1):
        parts.append(f"\n[Table {table_idx}]")
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))

    if ocr_images:
        try:
            with zipfile.ZipFile(BytesIO(file_bytes)) as z:
                media = [n for n in z.namelist() if n.startswith("word/media/")]
                for idx, name in enumerate(media, start=1):
                    suffix = Path(name).suffix.lower()
                    if suffix in IMAGE_EXTENSIONS:
                        parts.append(f"\n[Embedded Image {idx}: {Path(name).name}]")
                        parts.append(ocr_image_bytes(z.read(name)))
        except Exception as exc:
            parts.append(f"\n[Embedded image extraction failed: {exc}]")

    return "\n".join(parts)

def extract_pptx_bytes(file_bytes: bytes, ocr_images: bool = False) -> str:
    prs = Presentation(BytesIO(file_bytes))
    parts = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"\n[Slide {slide_idx}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    parts.append(text)

    if ocr_images:
        try:
            with zipfile.ZipFile(BytesIO(file_bytes)) as z:
                media = [n for n in z.namelist() if n.startswith("ppt/media/")]
                for idx, name in enumerate(media, start=1):
                    suffix = Path(name).suffix.lower()
                    if suffix in IMAGE_EXTENSIONS:
                        parts.append(f"\n[Embedded Image {idx}: {Path(name).name}]")
                        parts.append(ocr_image_bytes(z.read(name)))
        except Exception as exc:
            parts.append(f"\n[Embedded image extraction failed: {exc}]")

    return "\n".join(parts)

def extract_excel_bytes(file_bytes: bytes, suffix: str) -> Tuple[str, List[str]]:
    parts = []
    sheet_names = []

    if suffix == ".csv":
        try:
            df = pd.read_csv(BytesIO(file_bytes), dtype=str).fillna("")
            parts.append(df.head(200).to_csv(index=False))
            return "\n".join(parts), ["CSV"]
        except Exception as exc:
            return f"[CSV file detected but could not be read: {exc}]", []

    try:
        xls = pd.ExcelFile(BytesIO(file_bytes))
    except Exception as exc:
        return f"[Excel file detected but could not be read: {exc}]", []

    sheet_names = xls.sheet_names
    for sheet in sheet_names:
        try:
            df = pd.read_excel(BytesIO(file_bytes), sheet_name=sheet, dtype=str)
            df = df.fillna("")
            parts.append(f"\n[Sheet: {sheet}]")
            if df.empty:
                parts.append("[Empty sheet]")
            else:
                parts.append(df.head(120).to_csv(index=False))
        except Exception as exc:
            parts.append(f"\n[Sheet: {sheet}]")
            parts.append(f"[Could not read sheet: {exc}]")

    return "\n".join(parts), sheet_names

def extract_pdf_bytes(file_bytes: bytes, ocr_images: bool = False, max_ocr_pages: int = 5) -> str:
    try:
        import fitz
    except Exception as exc:
        return f"[PDF detected but PyMuPDF is unavailable: {exc}]"

    parts = []
    doc = fitz.open(stream=file_bytes, filetype="pdf")

    for page_idx, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        parts.append(f"\n[PDF Page {page_idx}]")
        if text:
            parts.append(text)
        else:
            parts.append("[No embedded text detected on this page.]")

        if ocr_images and page_idx <= max_ocr_pages and len(text) < 80:
            try:
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                img_bytes = pix.tobytes("png")
                parts.append("[OCR attempt]")
                parts.append(ocr_image_bytes(img_bytes))
            except Exception as exc:
                parts.append(f"[PDF OCR failed: {exc}]")

    return "\n".join(parts)

def convert_legacy_office_bytes(file_bytes: bytes, original_name: str, target_ext: str) -> Tuple[bytes, str]:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise RuntimeError("LibreOffice is not installed or not found in PATH. Cannot convert legacy Office file.")

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp = Path(tmpdir)
        input_path = tmp / original_name
        input_path.write_bytes(file_bytes)

        convert_to = "docx" if target_ext == ".docx" else "pptx"
        subprocess.run(
            [soffice, "--headless", "--convert-to", convert_to, "--outdir", str(tmp), str(input_path)],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )

        converted = input_path.with_suffix(target_ext)
        if not converted.exists():
            candidates = list(tmp.glob(f"*{target_ext}"))
            if not candidates:
                raise RuntimeError("LibreOffice conversion finished but no converted file was found.")
            converted = candidates[0]

        return converted.read_bytes(), converted.name

def extract_text_from_bytes(name: str, file_bytes: bytes, source_path: str, ocr_images: bool = False) -> Dict[str, Any]:
    suffix = Path(name).suffix.lower()
    warnings = []
    sheet_names = []

    try:
        if suffix == ".docx":
            text = extract_docx_bytes(file_bytes, ocr_images=ocr_images)

        elif suffix == ".doc":
            try:
                converted_bytes, converted_name = convert_legacy_office_bytes(file_bytes, Path(name).name, ".docx")
                text = extract_docx_bytes(converted_bytes, ocr_images=ocr_images)
                warnings.append(f"Legacy .doc converted through LibreOffice to {converted_name}.")
            except Exception as exc:
                text = f"[Legacy .doc file detected but conversion failed: {exc}]"
                warnings.append(str(exc))

        elif suffix in [".xlsx", ".xls", ".xlsm", ".csv"]:
            text, sheet_names = extract_excel_bytes(file_bytes, suffix)

        elif suffix == ".pptx":
            text = extract_pptx_bytes(file_bytes, ocr_images=ocr_images)

        elif suffix == ".ppt":
            try:
                converted_bytes, converted_name = convert_legacy_office_bytes(file_bytes, Path(name).name, ".pptx")
                text = extract_pptx_bytes(converted_bytes, ocr_images=ocr_images)
                warnings.append(f"Legacy .ppt converted through LibreOffice to {converted_name}.")
            except Exception as exc:
                text = f"[Legacy .ppt file detected but conversion failed: {exc}]"
                warnings.append(str(exc))

        elif suffix == ".pdf":
            text = extract_pdf_bytes(file_bytes, ocr_images=ocr_images)

        elif suffix in IMAGE_EXTENSIONS:
            text = ocr_image_bytes(file_bytes) if ocr_images else "[Image file detected. Enable OCR to extract text.]"

        elif suffix in [".txt", ".md"]:
            text = safe_decode(file_bytes)

        else:
            text = f"[Unsupported file type: {suffix}]"

    except Exception as exc:
        text = f"[Failed to extract text: {exc}]"
        warnings.append(str(exc))

    return {
        "document_name": Path(name).name,
        "source_path": source_path,
        "text": text,
        "file_type": suffix,
        "sheet_names": sheet_names,
        "extraction_warnings": warnings,
    }

def extract_text_from_upload(uploaded_file, ocr_images: bool = False) -> Dict[str, Any]:
    return extract_text_from_bytes(uploaded_file.name, uploaded_file.getvalue(), "uploaded_file", ocr_images=ocr_images)

def extract_docs_from_zip_upload(uploaded_file, ocr_images: bool = False) -> List[Dict[str, Any]]:
    docs = []
    file_bytes = uploaded_file.getvalue()

    with zipfile.ZipFile(BytesIO(file_bytes)) as z:
        for member in z.infolist():
            if member.is_dir():
                continue

            member_name = member.filename
            path = Path(member_name)
            suffix = path.suffix.lower()

            if suffix not in SUPPORTED_EXTENSIONS:
                continue
            if path.name.startswith("~$") or "__MACOSX" in path.parts:
                continue

            with z.open(member) as f:
                content = f.read()

            docs.append(
                extract_text_from_bytes(
                    name=member_name,
                    file_bytes=content,
                    source_path=f"{uploaded_file.name}/{member_name}",
                    ocr_images=ocr_images,
                )
            )

    return docs

# ============================================================
# AI Client
# ============================================================

def get_default_ai_config() -> Dict[str, Any]:
    return {
        "mode": os.getenv("AI_MODE", "heuristic"),
        "api_key": os.getenv("AI_API_KEY", ""),
        "base_url": os.getenv("AI_BASE_URL", ""),
        "model": os.getenv("AI_MODEL", ""),
    }

def safe_json_parse(text: str) -> Dict[str, Any]:
    if not text:
        raise ValueError("Empty AI response.")

    cleaned = text.strip()

    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```[a-zA-Z]*\n?", "", cleaned)
        cleaned = re.sub(r"\n?```$", "", cleaned)

    try:
        return json.loads(cleaned)
    except Exception:
        pass

    match = re.search(r"\{[\s\S]*\}", cleaned)
    if match:
        return json.loads(match.group(0))

    raise ValueError("Could not parse JSON from AI response.")

def normalize_result(data: Dict[str, Any], metadata: Dict[str, Any]) -> Dict[str, Any]:
    result = {}
    for field in REQUIRED_FIELDS:
        result[field] = data.get(field, "")

    result["document_name"] = result["document_name"] or metadata.get("document_name", "")
    result["source_path"] = result["source_path"] or metadata.get("source_path", "")
    result["file_type"] = result["file_type"] or metadata.get("file_type", "")

    if isinstance(result["needs_human_review"], str):
        result["needs_human_review"] = result["needs_human_review"].strip().lower() in ["true", "yes", "y", "1"]

    if result["risk_level"] not in ["Low", "Medium", "High"]:
        result["risk_level"] = "Medium"

    if result["can_enter_kb"] not in ["Yes", "No", "Needs Review"]:
        result["can_enter_kb"] = "Needs Review"

    return result

def heuristic_analyze(doc_text: str, metadata: Dict[str, Any]) -> Dict[str, Any]:
    name = metadata.get("document_name", "")
    lower = (name + "\n" + doc_text[:4000]).lower()

    document_type = "Unknown"
    department = "Unknown"
    business_process = "Unknown"
    knowledge_category = "General Business Knowledge"

    if "supplier" in lower or "vendor" in lower:
        department = "Procurement"
        business_process = "Supplier / Vendor Management"
    if "purchase" in lower or "procurement" in lower:
        department = "Procurement"
        business_process = "Purchase Approval"
    if "quality" in lower or "inspection" in lower or "complaint" in lower:
        department = "Quality"
        business_process = "Quality Management"
    if "maintenance" in lower or "equipment" in lower or "machine" in lower:
        department = "Maintenance"
        business_process = "Equipment Maintenance"
    if "production" in lower or "abnormal" in lower:
        department = "Production"
        business_process = "Production Issue Handling"
    if "safety" in lower or "ppe" in lower or "emergency" in lower:
        department = "EHS / Safety"
        business_process = "Safety Operation"
    if "training" in lower or "new employee" in lower or "onboarding" in lower:
        department = "HR / Operations"
        business_process = "Employee Onboarding"

    suffix = metadata.get("file_type", "")
    if "sop" in lower or "procedure" in lower or "process" in lower:
        document_type = "SOP / Process Document"
        knowledge_category = "Process Knowledge"
    elif suffix in [".xlsx", ".xls", ".xlsm", ".csv"]:
        document_type = "Spreadsheet / Operational Record"
        knowledge_category = "Operational Record / Checklist"
    elif suffix in [".pdf"]:
        document_type = "PDF Document"
        knowledge_category = "Documented Business Knowledge"
    elif suffix in [".ppt", ".pptx"]:
        document_type = "Presentation / Training Material"
        knowledge_category = "Training / Proposal Knowledge"
    elif suffix in IMAGE_EXTENSIONS:
        document_type = "Image / Screenshot"
        knowledge_category = "Visual Reference"
    elif "meeting" in lower or "action item" in lower:
        document_type = "Meeting Notes"
        knowledge_category = "Decision / Action Item Record"
    elif "manual" in lower or "guide" in lower:
        document_type = "Manual / Guide"
        knowledge_category = "Training / Reference Knowledge"

    risk_reasons = []
    if "version" not in lower and " v1" not in lower and " v2" not in lower and "_v1" not in lower and "_v2" not in lower:
        risk_reasons.append("Version is not clearly identified.")
    if "owner" not in lower and "responsible" not in lower and "approved by" not in lower:
        risk_reasons.append("Document owner or responsible role may be unclear.")
    if "unknown" in lower or "tbd" in lower or "to be confirmed" in lower or "needs review" in lower:
        risk_reasons.append("Contains unresolved or unclear information.")
    if "may" in lower or "should" in lower or "if needed" in lower or "as appropriate" in lower:
        risk_reasons.append("Some rules appear conditional or vague.")
    if "[failed" in lower or "[legacy" in lower or "ocr is unavailable" in lower:
        risk_reasons.append("Text extraction may be incomplete.")

    risk_level = "Low"
    if len(risk_reasons) >= 2:
        risk_level = "High"
    elif len(risk_reasons) == 1:
        risk_level = "Medium"

    summary = doc_text.strip().replace("\n", " ")[:450] or "No readable text extracted."
    can_enter_kb = "Yes" if risk_level == "Low" else "Needs Review"

    return normalize_result({
        "document_name": name,
        "source_path": metadata.get("source_path", ""),
        "file_type": metadata.get("file_type", ""),
        "document_type": document_type,
        "department": department,
        "business_process": business_process,
        "knowledge_category": knowledge_category,
        "summary": summary,
        "risk_level": risk_level,
        "risk_reason": " ".join(risk_reasons) if risk_reasons else "No major risk detected by fallback scan.",
        "recommended_action": "Review metadata and confirm whether this document is the current source of truth." if risk_level != "Low" else "Can be considered for knowledge base preparation.",
        "can_enter_kb": can_enter_kb,
        "needs_human_review": risk_level != "Low"
    }, metadata)

def call_gemini(prompt: str, config: Dict[str, Any]) -> str:
    api_key = config.get("api_key", "")
    model = config.get("model") or "gemini-2.5-flash"
    if not api_key:
        raise RuntimeError("API key is empty.")

    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
    payload = {
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {"temperature": 0.1}
    }

    response = requests.post(url, params={"key": api_key}, json=payload, timeout=90)
    response.raise_for_status()
    data = response.json()
    return data["candidates"][0]["content"]["parts"][0]["text"]

def call_openai_compatible(prompt: str, config: Dict[str, Any]) -> str:
    api_key = config.get("api_key", "")
    base_url = (config.get("base_url") or "").rstrip("/")
    model = config.get("model", "")

    if not api_key:
        raise RuntimeError("API key is empty.")
    if not base_url:
        raise RuntimeError("Base URL is empty.")
    if not model:
        raise RuntimeError("Model is empty.")

    url = base_url if base_url.endswith("/chat/completions") else f"{base_url}/chat/completions"
    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": "You are a precise enterprise knowledge inventory assistant. Follow output format strictly."},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.1
    }
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}

    response = requests.post(url, headers=headers, json=payload, timeout=90)
    response.raise_for_status()
    data = response.json()
    return data["choices"][0]["message"]["content"]

def call_anthropic(prompt: str, config: Dict[str, Any]) -> str:
    api_key = config.get("api_key", "")
    base_url = (config.get("base_url") or "https://api.anthropic.com").rstrip("/")
    model = config.get("model") or "claude-3-5-haiku-latest"

    if not api_key:
        raise RuntimeError("API key is empty.")

    url = f"{base_url}/v1/messages"
    payload = {
        "model": model,
        "max_tokens": 4096,
        "temperature": 0.1,
        "messages": [{"role": "user", "content": prompt}]
    }
    headers = {
        "x-api-key": api_key,
        "anthropic-version": "2023-06-01",
        "Content-Type": "application/json"
    }

    response = requests.post(url, headers=headers, json=payload, timeout=90)
    response.raise_for_status()
    data = response.json()
    return data["content"][0]["text"]

def call_llm(prompt: str, config: Dict[str, Any]) -> str:
    mode = config.get("mode", "heuristic")
    if mode == "gemini":
        return call_gemini(prompt, config)
    if mode == "anthropic":
        return call_anthropic(prompt, config)
    if mode == "openai_compatible":
        return call_openai_compatible(prompt, config)
    raise RuntimeError("No API mode selected.")

def analyze_document(doc_text: str, metadata: Dict[str, Any], max_chars: int, ai_config: Dict[str, Any]) -> Dict[str, Any]:
    if ai_config.get("mode") == "heuristic":
        return heuristic_analyze(doc_text, metadata)

    prompt = (
        CLASSIFY_PROMPT
        .replace("{{metadata_json}}", json.dumps(metadata, ensure_ascii=False, indent=2))
        .replace("{{document_text}}", doc_text[:max_chars])
    )

    try:
        raw = call_llm(prompt, ai_config)
        data = safe_json_parse(raw)
        return normalize_result(data, metadata)
    except Exception as exc:
        fallback = heuristic_analyze(doc_text, metadata)
        fallback["risk_level"] = "High"
        fallback["risk_reason"] = f"AI call failed or returned invalid JSON. Fallback used. Error: {exc}"
        fallback["recommended_action"] = "Check endpoint, API key, model, or output format, then rerun analysis."
        fallback["can_enter_kb"] = "Needs Review"
        fallback["needs_human_review"] = True
        return fallback

def heuristic_process_map(inventory: List[Dict[str, Any]]) -> str:
    groups = {}
    for item in inventory:
        dept = item.get("department") or "Unknown"
        proc = item.get("business_process") or "Unknown"
        groups.setdefault(dept, {}).setdefault(proc, []).append(item)

    lines = ["# Process Map", "", "Generated from the knowledge inventory.", ""]
    for dept, processes in groups.items():
        lines.append(f"## {dept}\n")
        for proc, docs in processes.items():
            lines.append(f"### {proc}\n")
            lines.append("Related documents:")
            for doc in docs:
                lines.append(f"- {doc.get('document_name')} — {doc.get('document_type')}")
            needs_review = [d for d in docs if d.get("needs_human_review")]
            lines.append(f"\nDocumentation status: {'Needs review' if needs_review else 'Generally usable'}")
            if needs_review:
                lines.append("Risk notes:")
                for doc in needs_review[:5]:
                    lines.append(f"- {doc.get('document_name')}: {doc.get('risk_reason')}")
            lines.append("")
    return "\n".join(lines)

def heuristic_risk_report(inventory: List[Dict[str, Any]]) -> str:
    high = [x for x in inventory if x.get("risk_level") == "High"]
    medium = [x for x in inventory if x.get("risk_level") == "Medium"]
    kb_ready = [x for x in inventory if x.get("can_enter_kb") == "Yes"]
    review = [x for x in inventory if x.get("needs_human_review")]

    lines = ["# Document Risk Report", "", "## Executive Summary", ""]
    lines.append(f"- Total documents analyzed: {len(inventory)}")
    lines.append(f"- High-risk documents: {len(high)}")
    lines.append(f"- Medium-risk documents: {len(medium)}")
    lines.append(f"- Documents needing human review: {len(review)}")
    lines.append(f"- Documents potentially ready for KB preparation: {len(kb_ready)}\n")

    lines.append("## Priority Review Items\n")
    if review:
        for doc in review:
            lines.append(f"### {doc.get('document_name')}")
            lines.append(f"- Department: {doc.get('department')}")
            lines.append(f"- Process: {doc.get('business_process')}")
            lines.append(f"- Risk level: {doc.get('risk_level')}")
            lines.append(f"- Risk reason: {doc.get('risk_reason')}")
            lines.append(f"- Recommended action: {doc.get('recommended_action')}\n")
    else:
        lines.append("No documents were flagged for human review.\n")

    lines.append("## Priority KB Candidates\n")
    if kb_ready:
        for doc in kb_ready:
            lines.append(f"- {doc.get('document_name')} — {doc.get('business_process')}")
    else:
        lines.append("No documents were marked as ready without review.")
    return "\n".join(lines)

def generate_process_map(inventory: List[Dict[str, Any]], ai_config: Dict[str, Any]) -> str:
    if ai_config.get("mode") == "heuristic":
        return heuristic_process_map(inventory)
    prompt = PROCESS_MAP_PROMPT.replace("{{inventory_json}}", json.dumps(inventory, ensure_ascii=False, indent=2))
    try:
        return call_llm(prompt, ai_config)
    except Exception as exc:
        return heuristic_process_map(inventory) + f"\n\n> AI generation failed; fallback used. Error: {exc}\n"

def generate_risk_report(inventory: List[Dict[str, Any]], ai_config: Dict[str, Any]) -> str:
    if ai_config.get("mode") == "heuristic":
        return heuristic_risk_report(inventory)
    prompt = RISK_REPORT_PROMPT.replace("{{inventory_json}}", json.dumps(inventory, ensure_ascii=False, indent=2))
    try:
        return call_llm(prompt, ai_config)
    except Exception as exc:
        return heuristic_risk_report(inventory) + f"\n\n> AI generation failed; fallback used. Error: {exc}\n"

# ============================================================
# Export
# ============================================================

def inventory_excel_bytes(df: pd.DataFrame) -> bytes:
    buffer = BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="Knowledge Inventory")
        ws = writer.book["Knowledge Inventory"]

        header_fill = PatternFill("solid", fgColor="D9EAF7")
        header_font = Font(bold=True)

        for cell in ws[1]:
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)

        for row in ws.iter_rows(min_row=2):
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)

        for col_idx, col in enumerate(ws.columns, start=1):
            max_len = 0
            for cell in col:
                value = "" if cell.value is None else str(cell.value)
                max_len = max(max_len, min(len(value), 60))
            ws.column_dimensions[get_column_letter(col_idx)].width = max(12, max_len + 2)

        ws.freeze_panes = "A2"

    return buffer.getvalue()

def markdown_bytes(text: str) -> bytes:
    return (text or "").encode("utf-8")

def report_zip_bytes(df: pd.DataFrame, process_map: str, risk_report: str) -> bytes:
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("knowledge_inventory.xlsx", inventory_excel_bytes(df))
        z.writestr("process_map.md", markdown_bytes(process_map))
        z.writestr("document_risk_report.md", markdown_bytes(risk_report))
    return buffer.getvalue()


# ============================================================
# Scope-aware Pre-correction
# ============================================================

def lower_risk_one_level(level: str) -> str:
    if level == "High":
        return "Medium"
    if level == "Medium":
        return "Low"
    return level or "Medium"

def is_actionable_answer(answer: str) -> bool:
    if not answer or not answer.strip():
        return False
    bad_tokens = ["unknown", "not sure", "unsure", "tbd", "to be confirmed", "不确定", "不知道", "待确认"]
    lower = answer.strip().lower()
    return not any(t in lower for t in bad_tokens)

def detect_issue_types(row: Dict[str, Any]) -> List[str]:
    reason = str(row.get("risk_reason", "")).lower()
    doc_type = str(row.get("document_type", "")).lower()
    file_type = str(row.get("file_type", "")).lower()
    summary = str(row.get("summary", "")).lower()

    issues = []

    if "version" in reason or "current" in reason or "latest" in reason:
        issues.append("version_unclear")
    if "owner" in reason or "responsible" in reason:
        issues.append("owner_unclear")
    if "unclear" in reason or "conditional" in reason or "vague" in reason:
        issues.append("rule_unclear")
    if "extraction" in reason or "ocr" in reason or "failed" in reason or "legacy" in reason:
        issues.append("extraction_incomplete")
    if "meeting" in doc_type or "presentation" in doc_type or file_type in [".ppt", ".pptx"]:
        issues.append("source_of_truth_unclear")
    if "conflict" in reason or "inconsistent" in reason or "conflicting" in summary:
        issues.append("conflict_detected")
    if str(row.get("risk_level", "")) == "High" and not issues:
        issues.append("high_risk_general")

    return list(dict.fromkeys(issues))

def build_issue_registry(df: pd.DataFrame) -> pd.DataFrame:
    rows = []
    if df is None or df.empty:
        return pd.DataFrame()

    for _, row in df.iterrows():
        row_dict = row.to_dict()
        issue_types = detect_issue_types(row_dict)

        for issue_type in issue_types:
            rows.append({
                "issue_id": f"ISS-{len(rows)+1:03d}",
                "document_name": row_dict.get("document_name", ""),
                "document_type": row_dict.get("document_type", ""),
                "department": row_dict.get("department", ""),
                "business_process": row_dict.get("business_process", ""),
                "risk_level": row_dict.get("risk_level", ""),
                "issue_type": issue_type,
                "uncertain_point": row_dict.get("risk_reason", ""),
                "impact_level": "High" if row_dict.get("risk_level") == "High" else "Medium",
                "confidence": "Low" if row_dict.get("risk_level") == "High" else "Medium",
                "evidence_excerpt": str(row_dict.get("summary", ""))[:350],
            })

    return pd.DataFrame(rows)

def issue_label(issue_type: str) -> str:
    labels = {
        "version_unclear": "Version / validity unclear",
        "owner_unclear": "Owner unclear",
        "rule_unclear": "Business rule unclear",
        "source_of_truth_unclear": "Source-of-truth unclear",
        "extraction_incomplete": "Extraction quality issue",
        "conflict_detected": "Potential conflict",
        "high_risk_general": "High-risk item",
    }
    return labels.get(issue_type, issue_type)

def build_answer_options(issue_type: str, scope_type: str, docs: pd.DataFrame) -> tuple[list[dict], str, str]:
    """Return Claude-like answer options plus a recommended choice."""
    count = len(docs["document_name"].dropna().unique().tolist()) if docs is not None and not docs.empty else 0

    if issue_type == "version_unclear":
        if scope_type == "Document-level":
            options = [
                {"label": "Treat this document as the current valid version", "effect": "lower_risk", "description": "Use when the business owner confirms this is the active approved file."},
                {"label": "Keep this document as Needs Review until version is confirmed", "effect": "keep_review", "description": "Safest choice when the version or approval status is still unclear."},
                {"label": "Treat this document as supporting material only", "effect": "supporting_only", "description": "Use when the file is useful context but should not become source of truth."},
                {"label": "Other / custom answer", "effect": "custom", "description": "Write your own correction rule."},
            ]
        else:
            options = [
                {"label": "Keep documents without clear version as Needs Review", "effect": "keep_review", "description": f"Safest batch rule for {count} affected document(s)."},
                {"label": "Treat the affected documents as current valid versions", "effect": "lower_risk", "description": "Use only if the upload batch is already confirmed as the current approved set."},
                {"label": "Treat the affected documents as supporting material only", "effect": "supporting_only", "description": "Useful when these files provide context but should not drive final AI answers."},
                {"label": "Other / custom answer", "effect": "custom", "description": "Write your own correction rule."},
            ]
        recommended = "Keep documents without clear version as Needs Review" if scope_type != "Document-level" else "Keep this document as Needs Review until version is confirmed"
        reason = "Recommended because version validity directly affects source-of-truth quality."

    elif issue_type == "owner_unclear":
        dept_candidates = [x for x in docs.get("department", pd.Series(dtype=str)).dropna().unique().tolist() if str(x).strip() and str(x).strip().lower() != "unknown"] if docs is not None and not docs.empty else []
        default_owner = dept_candidates[0] if len(dept_candidates) == 1 else "the selected business owner"
        options = [
            {"label": f"Assign {default_owner} as the owner for this scope", "effect": "lower_risk", "description": "Use when the selected scope clearly belongs to one responsible team."},
            {"label": "Keep owner as Needs Review", "effect": "keep_review", "description": "Use when ownership still needs confirmation."},
            {"label": "Apply owner only to selected documents, not future matches", "effect": "lower_risk", "description": "Use when the answer is valid only for this correction round."},
            {"label": "Other / custom answer", "effect": "custom", "description": "Write the exact owner or ownership rule."},
        ]
        recommended = f"Assign {default_owner} as the owner for this scope" if default_owner != "the selected business owner" else "Keep owner as Needs Review"
        reason = "Recommended based on whether this scope has a clear shared department."

    elif issue_type == "source_of_truth_unclear":
        options = [
            {"label": "Treat these documents as supporting material, not source of truth", "effect": "supporting_only", "description": "Best for meeting notes, presentations, drafts, and informal records."},
            {"label": "Treat these documents as source of truth", "effect": "lower_risk", "description": "Use only if these files are approved authoritative references."},
            {"label": "Keep as Needs Review until source-of-truth status is approved", "effect": "keep_review", "description": "Safe option when approval status is unclear."},
            {"label": "Other / custom answer", "effect": "custom", "description": "Write your own source-of-truth rule."},
        ]
        recommended = "Treat these documents as supporting material, not source of truth"
        reason = "Recommended because non-SOP documents often need approval before entering a knowledge base."

    elif issue_type == "rule_unclear":
        options = [
            {"label": "Keep unclear rules as Needs Review until the business rule is confirmed", "effect": "keep_review", "description": "Safest when the rule affects process guidance or AI answers."},
            {"label": "Apply one shared business rule to this scope", "effect": "lower_risk", "description": "Use when you can provide the confirmed rule in the custom field."},
            {"label": "Mark as not KB-ready because the rule is too ambiguous", "effect": "not_kb_ready", "description": "Use when ambiguity is high and cannot be resolved in this round."},
            {"label": "Other / custom answer", "effect": "custom", "description": "Write the exact clarified rule."},
        ]
        recommended = "Keep unclear rules as Needs Review until the business rule is confirmed"
        reason = "Recommended because unclear business rules can create wrong operational answers."

    elif issue_type == "extraction_incomplete":
        options = [
            {"label": "Keep as Needs Review due to extraction quality", "effect": "keep_review", "description": "Safest when OCR or conversion may have missed content."},
            {"label": "Use extracted text for preliminary analysis only", "effect": "supporting_only", "description": "Allows rough inventory while blocking final knowledge ingestion."},
            {"label": "Exclude from KB until the file is re-uploaded or OCR is fixed", "effect": "not_kb_ready", "description": "Use when extraction quality is not acceptable."},
            {"label": "Other / custom answer", "effect": "custom", "description": "Write your own handling rule."},
        ]
        recommended = "Keep as Needs Review due to extraction quality"
        reason = "Recommended because incomplete extraction can hide important clauses or steps."

    elif issue_type == "conflict_detected":
        options = [
            {"label": "Keep as Needs Review until the conflict is resolved", "effect": "keep_review", "description": "Safest when multiple documents may disagree."},
            {"label": "Use the most recent approved document as source of truth", "effect": "lower_risk", "description": "Use only when approval and version information are reliable."},
            {"label": "Ask the process owner to confirm the current standard", "effect": "keep_review", "description": "Use when ownership is known but the rule is not."},
            {"label": "Other / custom answer", "effect": "custom", "description": "Write the specific rule or source to use."},
        ]
        recommended = "Keep as Needs Review until the conflict is resolved"
        reason = "Recommended because conflicts should not be published as trusted knowledge."

    else:
        options = [
            {"label": "Keep these high-risk items as Needs Review", "effect": "keep_review", "description": "Safe default for uncertain or incomplete documents."},
            {"label": "Treat them as low-confidence supporting material", "effect": "supporting_only", "description": "Use when content is useful but should not drive final answers."},
            {"label": "Apply a shared correction rule to this scope", "effect": "lower_risk", "description": "Use when you can provide the correction rule in the custom field."},
            {"label": "Other / custom answer", "effect": "custom", "description": "Write your own correction rule."},
        ]
        recommended = "Keep these high-risk items as Needs Review"
        reason = "Recommended because high-risk documents require a conservative default."

    return options, recommended, reason


def build_question(issue_type, docs, scope_type, context_label):
    affected = docs["document_name"].dropna().unique().tolist()
    count = len(affected)

    base = {
        "question_group_id": "",
        "scope_type": scope_type,
        "issue_type": issue_type,
        "issue_label": issue_label(issue_type),
        "affected_documents": affected,
        "affected_document_count": count,
        "context_label": context_label,
        "user_answer": "",
        "apply_scope": "Affected documents in this question",
        "status": "Open",
    }

    if issue_type == "version_unclear":
        if scope_type == "Document-level":
            q = "Should this document be treated as the current valid version?"
            assump = "This document may not be the latest approved version."
        else:
            q = f"I found {count} document(s) with unclear version or validity. How should these documents be treated?"
            assump = "Documents without clear version information should remain Needs Review until confirmed."
        why = "Version validity affects whether the document can be used as a trusted source for knowledge ingestion."

    elif issue_type == "owner_unclear":
        if scope_type == "Document-level":
            q = "Who is the owner or responsible team for this document?"
            assump = "The responsible owner is not clearly stated."
        else:
            q = f"I found {count} document(s) with unclear owner information. Is there a default owner or responsible team for this scope?"
            assump = "A shared owner may apply to this group of documents."
        why = "Ownership matters because business users need to know who can approve, update, or validate the knowledge."

    elif issue_type == "source_of_truth_unclear":
        if scope_type == "Document-level":
            q = "Should this document be treated as source of truth, or only as supporting material?"
            assump = "This document may be supporting material rather than an authoritative source."
        else:
            q = f"I found {count} document(s) that may be supporting materials, such as meeting notes or presentations. Should they be treated as source of truth?"
            assump = "Meeting notes and presentations should generally be supporting materials unless explicitly approved."
        why = "Source-of-truth status affects whether content can be directly used for RAG or enterprise assistant answers."

    elif issue_type == "rule_unclear":
        if scope_type == "Document-level":
            q = "What is the correct business rule or interpretation for the unclear part in this document?"
            assump = "Some rule wording is conditional, vague, or incomplete."
        else:
            q = f"I found {count} document(s) with vague or conditional rule wording. Is there a shared rule that should be applied to this scope?"
            assump = "The current wording may need human clarification before knowledge ingestion."
        why = "Unclear rules can lead to unreliable AI answers or inconsistent process guidance."

    elif issue_type == "extraction_incomplete":
        if scope_type == "Document-level":
            q = "Can this document be analyzed based on the extracted text, or should it remain Needs Review due to extraction limitations?"
            assump = "The extracted text may be incomplete."
        else:
            q = f"I found {count} document(s) with extraction or OCR limitations. Should these remain Needs Review until manually checked?"
            assump = "Incomplete extraction should block direct knowledge ingestion."
        why = "If text extraction is incomplete, AI analysis may miss important content."

    elif issue_type == "conflict_detected":
        q = f"I found potential conflicts in {count} document(s). Which rule or source should be treated as the current standard?"
        assump = "Conflicting information should not be published without confirmation."
        why = "Conflicts must be resolved before the knowledge can be trusted."

    else:
        q = f"I found {count} high-risk document(s) in this scope. Is there a shared rule or clarification that can reduce the review risk?"
        assump = "High-risk documents should remain Needs Review unless clarified."
        why = "High-risk items may affect downstream AI answer reliability."

    options, recommended, recommended_reason = build_answer_options(issue_type, scope_type, docs)

    base.update({
        "why_it_matters": why,
        "ai_assumption": assump,
        "question_to_user": q,
        "answer_options": options,
        "recommended_option": recommended,
        "recommended_reason": recommended_reason,
    })
    return base

def generate_scope_questions(df: pd.DataFrame, scope_type: str, scope_df: pd.DataFrame) -> pd.DataFrame:
    issue_df = build_issue_registry(scope_df)
    if issue_df.empty:
        return pd.DataFrame()

    max_questions = {
        "Batch-level": 8,
        "Segment-level": 6,
        "Selected-documents": 5,
        "Document-level": 5,
    }.get(scope_type, 6)

    questions = []

    if scope_type in ["Batch-level", "Segment-level", "Selected-documents"]:
        priority = {
            "conflict_detected": 0,
            "version_unclear": 1,
            "owner_unclear": 2,
            "source_of_truth_unclear": 3,
            "rule_unclear": 4,
            "extraction_incomplete": 5,
            "high_risk_general": 6,
        }

        grouped = []
        for issue_type, g in issue_df.groupby("issue_type"):
            affected_docs = g["document_name"].dropna().unique().tolist()
            high_count = int((g["risk_level"] == "High").sum())
            grouped.append({
                "issue_type": issue_type,
                "g": g,
                "affected_count": len(affected_docs),
                "high_count": high_count,
                "priority": priority.get(issue_type, 99),
            })

        grouped = sorted(grouped, key=lambda x: (x["priority"], -x["high_count"], -x["affected_count"]))

        context_label = "entire batch" if scope_type == "Batch-level" else "selected segment"
        for item in grouped[:max_questions]:
            q = build_question(item["issue_type"], item["g"], scope_type, context_label)
            questions.append(q)

    else:
        doc_name = scope_df["document_name"].iloc[0] if not scope_df.empty else "selected document"
        context_label = doc_name

        for issue_type, g in issue_df.groupby("issue_type"):
            q = build_question(issue_type, g, scope_type, context_label)
            questions.append(q)
            if len(questions) >= max_questions:
                break

    for i, q in enumerate(questions, start=1):
        q["question_group_id"] = f"QG-{i:03d}"

    return pd.DataFrame(questions)

def infer_answer_effect(answer: str) -> str:
    lower = (answer or "").lower()
    if not lower.strip():
        return "skip"
    if "other / custom" in lower and "user clarification:" not in lower:
        return "skip"
    if "not kb-ready" in lower or "exclude from kb" in lower or "too ambiguous" in lower:
        return "not_kb_ready"
    if "supporting material" in lower or "preliminary analysis only" in lower or "not source of truth" in lower:
        return "supporting_only"
    if "needs review" in lower or "until" in lower or "ask the process owner" in lower:
        return "keep_review"
    if is_actionable_answer(answer):
        return "lower_risk"
    return "keep_review"


def calculate_corrected_status(before_risk: str, before_kb: str, answer: str) -> tuple[str, str, bool]:
    effect = infer_answer_effect(answer)
    if effect == "skip":
        return before_risk, before_kb, True
    if effect == "keep_review":
        return before_risk, "Needs Review", True
    if effect == "supporting_only":
        after_risk = lower_risk_one_level(before_risk)
        return after_risk, "Needs Review", True
    if effect == "not_kb_ready":
        return "High", "No", True
    if effect == "lower_risk":
        after_risk = lower_risk_one_level(before_risk)
        after_kb = "Yes" if after_risk == "Low" else "Needs Review"
        after_review = False if after_risk == "Low" else True
        return after_risk, after_kb, after_review
    return before_risk, "Needs Review", True


def apply_precorrections(
    inventory_df: pd.DataFrame,
    questions_df: pd.DataFrame,
    answers: Dict[str, str],
    apply_scopes: Dict[str, str],
    selected_scope_docs: List[str],
) -> tuple[pd.DataFrame, pd.DataFrame]:
    corrected = inventory_df.copy()

    if "corrected_risk_level" not in corrected.columns:
        corrected["corrected_risk_level"] = corrected["risk_level"]
    if "corrected_can_enter_kb" not in corrected.columns:
        corrected["corrected_can_enter_kb"] = corrected["can_enter_kb"]
    if "corrected_needs_human_review" not in corrected.columns:
        corrected["corrected_needs_human_review"] = corrected["needs_human_review"]
    if "correction_notes" not in corrected.columns:
        corrected["correction_notes"] = ""

    ledger_rows = []

    if questions_df is None or questions_df.empty:
        return corrected, pd.DataFrame()

    for _, q in questions_df.iterrows():
        qid = q["question_group_id"]
        answer = answers.get(qid, "").strip()
        if not answer:
            continue

        apply_scope = apply_scopes.get(qid, "Affected documents in this question")
        affected_docs = q.get("affected_documents", [])
        if isinstance(affected_docs, str):
            try:
                affected_docs = json.loads(affected_docs)
            except Exception:
                affected_docs = [affected_docs]

        issue_type = q.get("issue_type", "")

        if apply_scope == "Current selected scope":
            target_docs = selected_scope_docs
        elif apply_scope == "Entire batch matching this issue":
            target_docs = []
            for _, row in corrected.iterrows():
                if issue_type in detect_issue_types(row.to_dict()):
                    target_docs.append(row.get("document_name", ""))
        elif apply_scope == "This document only":
            target_docs = affected_docs[:1]
        else:
            target_docs = affected_docs

        target_docs = list(dict.fromkeys([d for d in target_docs if d]))

        for doc_name in target_docs:
            mask = corrected["document_name"] == doc_name
            if not mask.any():
                continue

            before_risk = corrected.loc[mask, "corrected_risk_level"].iloc[0]
            before_kb = corrected.loc[mask, "corrected_can_enter_kb"].iloc[0]

            after_risk, after_kb, after_review = calculate_corrected_status(before_risk, before_kb, answer)

            note = f"[{qid} | {issue_label(issue_type)} | {apply_scope}] {answer}"
            existing = str(corrected.loc[mask, "correction_notes"].iloc[0] or "")
            corrected.loc[mask, "correction_notes"] = (existing + "\n" + note).strip()
            corrected.loc[mask, "corrected_risk_level"] = after_risk
            corrected.loc[mask, "corrected_can_enter_kb"] = after_kb
            corrected.loc[mask, "corrected_needs_human_review"] = after_review

            ledger_rows.append({
                "correction_id": f"COR-{len(ledger_rows)+1:03d}",
                "question_group_id": qid,
                "scope_type": q.get("scope_type", ""),
                "issue_type": issue_type,
                "document_name": doc_name,
                "user_answer": answer,
                "apply_scope": apply_scope,
                "before_risk_level": before_risk,
                "after_risk_level": after_risk,
                "before_can_enter_kb": before_kb,
                "after_can_enter_kb": after_kb,
                "correction_note": note,
                "status": "Applied",
            })

    return corrected, pd.DataFrame(ledger_rows)

def corrected_view_for_report(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty or "corrected_risk_level" not in df.columns:
        return df

    report_df = df.copy()
    report_df["risk_level_original"] = report_df["risk_level"]
    report_df["can_enter_kb_original"] = report_df["can_enter_kb"]
    report_df["needs_human_review_original"] = report_df["needs_human_review"]

    report_df["risk_level"] = report_df["corrected_risk_level"]
    report_df["can_enter_kb"] = report_df["corrected_can_enter_kb"]
    report_df["needs_human_review"] = report_df["corrected_needs_human_review"]
    return report_df

def dataframe_excel_bytes(df: pd.DataFrame, sheet_name: str = "Sheet1") -> bytes:
    buffer = BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name=sheet_name[:31])
        ws = writer.book[sheet_name[:31]]

        header_fill = PatternFill("solid", fgColor="D9EAF7")
        header_font = Font(bold=True)

        for cell in ws[1]:
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)

        for row in ws.iter_rows(min_row=2):
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)

        for col_idx, col in enumerate(ws.columns, start=1):
            max_len = 0
            for cell in col:
                value = "" if cell.value is None else str(cell.value)
                max_len = max(max_len, min(len(value), 60))
            ws.column_dimensions[get_column_letter(col_idx)].width = max(12, max_len + 2)

        ws.freeze_panes = "A2"

    return buffer.getvalue()

def report_zip_bytes_with_ledger(
    df: pd.DataFrame,
    process_map: str,
    risk_report: str,
    ledger_df: pd.DataFrame | None = None,
    corrected_docs_zip: bytes | None = None,
) -> bytes:
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("knowledge_inventory_corrected.xlsx", inventory_excel_bytes(df))
        z.writestr("process_map_corrected.md", markdown_bytes(process_map))
        z.writestr("document_risk_report_corrected.md", markdown_bytes(risk_report))
        if ledger_df is not None and not ledger_df.empty:
            z.writestr("correction_ledger.xlsx", dataframe_excel_bytes(ledger_df, "Correction Ledger"))
        if corrected_docs_zip:
            # Keep corrected documents as a separate package inside the full bundle.
            z.writestr("corrected_documents.zip", corrected_docs_zip)
    return buffer.getvalue()


def sanitize_filename(name: str) -> str:
    safe = re.sub(r"[^A-Za-z0-9._-]+", "_", str(name or "document")).strip("_")
    return safe[:120] or "document"


def get_raw_doc_text(raw_documents: Dict[str, Any], document_name: str) -> str:
    raw = raw_documents.get(document_name, {}) if isinstance(raw_documents, dict) else {}
    text = raw.get("text", "") if isinstance(raw, dict) else ""
    return text or "[Original extracted text is not available in session. This draft is generated from the inventory record and correction ledger.]"


def rows_for_doc(df: pd.DataFrame, doc_name: str) -> Dict[str, Any]:
    if df is None or df.empty or "document_name" not in df.columns:
        return {}
    match = df[df["document_name"] == doc_name]
    if match.empty:
        return {}
    return match.iloc[0].to_dict()


def build_fallback_corrected_draft(
    row: Dict[str, Any],
    raw_text: str,
    corrections: List[Dict[str, Any]],
) -> str:
    doc_name = row.get("document_name", "Unknown document")
    corrected_risk = row.get("risk_level", row.get("corrected_risk_level", ""))
    corrected_kb = row.get("can_enter_kb", row.get("corrected_can_enter_kb", ""))
    summary = row.get("summary", "")
    recommendation = row.get("recommended_action", "")

    lines = [
        f"# Corrected Draft — {doc_name}",
        "",
        "> This document was generated after guided pre-correction. The original source file is not overwritten.",
        "",
        "## Corrected Document Status",
        "",
        f"- Document type: {row.get('document_type', 'Unknown')}",
        f"- Department: {row.get('department', 'Unknown')}",
        f"- Business process: {row.get('business_process', 'Unknown')}",
        f"- Knowledge category: {row.get('knowledge_category', 'Unknown')}",
        f"- Corrected risk level: {corrected_risk}",
        f"- Corrected KB readiness: {corrected_kb}",
        "",
        "## Applied Corrections",
        "",
    ]

    if corrections:
        for c in corrections:
            lines.extend([
                f"### {c.get('correction_id', '')} — {c.get('issue_type', '')}",
                f"- User clarification: {c.get('user_answer', '')}",
                f"- Apply scope: {c.get('apply_scope', '')}",
                f"- Risk level: {c.get('before_risk_level', '')} → {c.get('after_risk_level', '')}",
                f"- KB readiness: {c.get('before_can_enter_kb', '')} → {c.get('after_can_enter_kb', '')}",
                "",
            ])
    else:
        lines.append("No correction ledger entries were found for this document.\n")

    lines.extend([
        "## Revised Working Summary",
        "",
        summary or "No summary available.",
        "",
        "## Recommended Next Action",
        "",
        recommendation or "Review and approve before publishing into a knowledge base.",
        "",
        "## Corrected Knowledge Draft",
        "",
        "The following draft preserves the extracted source content and adds the confirmed pre-correction decisions above. In a production deployment, this section can be replaced by a model-rewritten SOP / FAQ / knowledge article based on the same correction ledger.",
        "",
        "### Source Content for Revision",
        "",
        raw_text[:18000] if raw_text else "[No extracted source content available.]",
        "",
        "---",
        "Generated by AI Knowledge Inventory — Guided Pre-correction.",
    ])

    return "\n".join(lines)


def generate_corrected_draft_markdown(
    row: Dict[str, Any],
    raw_text: str,
    corrections: List[Dict[str, Any]],
    ai_config: Dict[str, Any],
) -> str:
    # In No API mode, generate a structured corrected draft deterministically.
    if ai_config.get("mode") == "heuristic":
        return build_fallback_corrected_draft(row, raw_text, corrections)

    prompt = f"""
You are an enterprise documentation specialist.

Create a corrected, business-readable Markdown draft based on the original extracted document text and the correction ledger.

Important rules:
- Do not invent unsupported business rules.
- Keep the original document meaning where possible.
- Apply only the user-confirmed corrections.
- If the source text is incomplete, clearly mark the affected section as Needs Review.
- Do not claim this is formally approved unless the correction says so.
- Output Markdown only.

Document inventory record:
{json.dumps(row, ensure_ascii=False, indent=2)}

Correction ledger entries for this document:
{json.dumps(corrections, ensure_ascii=False, indent=2)}

Original extracted document text:
{raw_text[:14000]}
""".strip()

    try:
        return call_llm(prompt, ai_config).strip()
    except Exception as exc:
        fallback = build_fallback_corrected_draft(row, raw_text, corrections)
        return fallback + f"\n\n> AI draft rewrite failed; structured fallback used. Error: {exc}\n"


def markdown_to_docx_bytes(markdown_text: str, title: str = "Corrected Draft") -> bytes:
    doc = Document()
    doc.add_heading(title, level=0)

    for raw_line in markdown_text.splitlines():
        line = raw_line.rstrip()
        if not line:
            doc.add_paragraph("")
            continue
        if line.startswith("# "):
            doc.add_heading(line[2:].strip(), level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:].strip(), level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:].strip(), level=3)
        elif line.startswith("- "):
            doc.add_paragraph(line[2:].strip(), style="List Bullet")
        elif line.startswith("> "):
            p = doc.add_paragraph(line[2:].strip())
            p.style = "Intense Quote"
        elif line.strip() == "---":
            doc.add_paragraph("────────────────────────")
        else:
            doc.add_paragraph(line)

    buffer = BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def corrected_documents_zip_bytes(
    corrected_report_df: pd.DataFrame,
    ledger_df: pd.DataFrame | None,
    raw_documents: Dict[str, Any],
    ai_config: Dict[str, Any],
) -> tuple[bytes, int]:
    buffer = BytesIO()

    if ledger_df is None or ledger_df.empty:
        return buffer.getvalue(), 0

    doc_names = ledger_df["document_name"].dropna().unique().tolist()
    doc_names = [d for d in doc_names if str(d).strip()]

    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as z:
        manifest_rows = []
        for doc_name in doc_names:
            row = rows_for_doc(corrected_report_df, doc_name)
            raw_text = get_raw_doc_text(raw_documents, doc_name)
            corrections = ledger_df[ledger_df["document_name"] == doc_name].to_dict(orient="records")
            md = generate_corrected_draft_markdown(row, raw_text, corrections, ai_config)

            base_name = sanitize_filename(Path(str(doc_name)).stem)
            md_name = f"corrected_documents/{base_name}_corrected.md"
            docx_name = f"corrected_documents/{base_name}_corrected.docx"
            z.writestr(md_name, md.encode("utf-8"))
            z.writestr(docx_name, markdown_to_docx_bytes(md, f"Corrected Draft — {doc_name}"))

            manifest_rows.append({
                "document_name": doc_name,
                "markdown_file": md_name,
                "word_file": docx_name,
                "correction_count": len(corrections),
            })

        if manifest_rows:
            manifest = pd.DataFrame(manifest_rows)
            z.writestr("corrected_documents_manifest.xlsx", dataframe_excel_bytes(manifest, "Corrected Docs"))

    return buffer.getvalue(), len(doc_names)

def render_precorrection_workbench(base_df: pd.DataFrame, ai_config: Dict[str, Any]):
    st.subheader("Guided Pre-correction")
    st.caption("Choose the scope, answer a small number of targeted questions, and decide exactly how broadly each answer should be applied.")

    st.markdown('<div class="guided-card">', unsafe_allow_html=True)
    c1, c2, c3, c4 = st.columns(4)
    with c1:
        scope_choice = st.selectbox(
            "Correction scope",
            ["Batch-level", "Segment-level", "Selected-documents", "Document-level"],
            help="Batch-level fixes shared issues. Segment-level focuses on a filtered group. Document-level fixes one file only."
        )

    scope_df = base_df.copy()

    with c2:
        if scope_choice == "Segment-level":
            dept_options = sorted(scope_df["department"].dropna().unique().tolist()) if "department" in scope_df else []
            selected_depts = st.multiselect("Department", dept_options, default=dept_options[:1] if dept_options else [])
            if selected_depts:
                scope_df = scope_df[scope_df["department"].isin(selected_depts)]
        else:
            st.write("")

    with c3:
        if scope_choice == "Segment-level":
            process_options = sorted(scope_df["business_process"].dropna().unique().tolist()) if "business_process" in scope_df else []
            selected_processes = st.multiselect("Process", process_options, default=[])
            if selected_processes:
                scope_df = scope_df[scope_df["business_process"].isin(selected_processes)]
        else:
            st.write("")

    with c4:
        if scope_choice == "Segment-level":
            risk_options = sorted(scope_df["risk_level"].dropna().unique().tolist()) if "risk_level" in scope_df else []
            selected_risks = st.multiselect("Risk", risk_options, default=risk_options)
            if selected_risks:
                scope_df = scope_df[scope_df["risk_level"].isin(selected_risks)]
        else:
            st.write("")

    if scope_choice == "Selected-documents":
        doc_options = base_df["document_name"].dropna().tolist()
        selected_docs = st.multiselect("Select documents to refine", doc_options)
        scope_df = base_df[base_df["document_name"].isin(selected_docs)] if selected_docs else base_df.iloc[0:0]

    elif scope_choice == "Document-level":
        doc_options = base_df["document_name"].dropna().tolist()
        one_doc = st.selectbox("Select one document", doc_options) if doc_options else None
        scope_df = base_df[base_df["document_name"] == one_doc] if one_doc else base_df.iloc[0:0]

    st.info(f"Current correction scope includes {len(scope_df)} document(s).")

    metric_cols = st.columns(4)
    metric_cols[0].metric("Scope documents", len(scope_df))
    if not scope_df.empty and "risk_level" in scope_df:
        metric_cols[1].metric("High risk in scope", int((scope_df["risk_level"] == "High").sum()))
        metric_cols[2].metric("Need review in scope", int(scope_df["needs_human_review"].astype(bool).sum()) if "needs_human_review" in scope_df else 0)
    issue_df = build_issue_registry(scope_df)
    metric_cols[3].metric("Detected issues", len(issue_df))

    gen_label = "Generate guided questions for this scope"
    if st.button(gen_label, type="primary", use_container_width=True):
        if scope_df.empty:
            st.warning("No documents in current scope.")
        else:
            qdf = generate_scope_questions(base_df, scope_choice, scope_df)
            st.session_state.correction_questions = qdf
            st.session_state.correction_scope_docs = scope_df["document_name"].dropna().tolist()
            st.session_state.correction_scope_type = scope_choice
    st.markdown('</div>', unsafe_allow_html=True)

    qdf = st.session_state.get("correction_questions")

    if qdf is not None and not qdf.empty:
        st.markdown("### Clarification Questions")
        st.caption(f"{len(qdf)} grouped question(s) generated. Pick an option, review the recommended choice, or write your own clarification.")

        answers = {}
        apply_scopes = {}

        for _, q in qdf.iterrows():
            qid = q["question_group_id"]
            with st.expander(f"{qid} · {q['issue_label']} · {q['affected_document_count']} affected document(s)", expanded=True):
                st.markdown(f"**Question**  \n{q['question_to_user']}")
                st.markdown(f"**Why it matters**  \n{q['why_it_matters']}")
                st.markdown(f"**AI assumption**  \n{q['ai_assumption']}")

                recommended = q.get("recommended_option", "")
                recommended_reason = q.get("recommended_reason", "")
                if recommended:
                    st.success(f"Recommended option: {recommended}\n\nReason: {recommended_reason}")

                docs = q.get("affected_documents", [])
                if isinstance(docs, str):
                    try:
                        docs = json.loads(docs)
                    except Exception:
                        docs = [docs]
                if docs:
                    st.caption("Affected documents: " + ", ".join(docs[:8]) + (" ..." if len(docs) > 8 else ""))

                raw_options = q.get("answer_options", [])
                if not isinstance(raw_options, list) or not raw_options:
                    raw_options = [{"label": "Keep as Needs Review", "description": "Safe default."}, {"label": "Other / custom answer", "description": "Write your own correction rule."}]

                labels = [opt.get("label", str(opt)) if isinstance(opt, dict) else str(opt) for opt in raw_options]
                recommended_idx = labels.index(recommended) if recommended in labels else 0
                display_labels = [f"⭐ Recommended — {x}" if x == recommended else x for x in labels]

                selected_display = st.radio(
                    "Choose one answer",
                    display_labels,
                    index=recommended_idx,
                    key=f"option_{qid}",
                    help="You can accept the recommended option or choose another option."
                )
                selected_label = selected_display.replace("⭐ Recommended — ", "", 1)

                selected_desc = ""
                for opt in raw_options:
                    if isinstance(opt, dict) and opt.get("label") == selected_label:
                        selected_desc = opt.get("description", "")
                        break
                if selected_desc:
                    st.caption(f"Option meaning: {selected_desc}")

                custom_note = st.text_area(
                    "Optional custom clarification / override",
                    key=f"custom_{qid}",
                    placeholder="Example: Owner should be Procurement Operations only for Supplier Onboarding documents created after 2025.",
                    height=85,
                )

                if selected_label == "Other / custom answer" and not custom_note.strip():
                    answers[qid] = ""
                else:
                    answers[qid] = selected_label if not custom_note.strip() else f"{selected_label}. User clarification: {custom_note.strip()}"

                scope_options = ["Affected documents in this question", "Current selected scope", "Entire batch matching this issue"]
                if q["affected_document_count"] >= 1:
                    scope_options.append("This document only")

                default_scope_idx = 0
                if q.get("scope_type") == "Document-level" and "This document only" in scope_options:
                    default_scope_idx = scope_options.index("This document only")
                elif q.get("scope_type") == "Batch-level":
                    default_scope_idx = scope_options.index("Affected documents in this question")

                apply_scopes[qid] = st.selectbox(
                    "Apply this answer to",
                    scope_options,
                    index=default_scope_idx,
                    key=f"apply_scope_{qid}",
                    help="This controls whether the correction affects one document, selected documents, or all matching documents."
                )

        if st.button("Apply Answers & Regenerate Corrected Report", type="primary", use_container_width=True):
            missing_custom = [qid for qid, ans in answers.items() if not ans.strip()]
            if missing_custom:
                st.warning("Some questions use 'Other / custom answer' but have no custom clarification. Please fill them in or choose another option.")
            else:
                corrected_df, ledger_df = apply_precorrections(
                    base_df,
                    qdf,
                    answers,
                    apply_scopes,
                    st.session_state.get("correction_scope_docs", []),
                )

                report_df = corrected_view_for_report(corrected_df)
                records = report_df.to_dict(orient="records")
                corrected_process_map = generate_process_map(records, ai_config=ai_config)
                corrected_risk_report = generate_risk_report(records, ai_config=ai_config)

                st.session_state.corrected_inventory = corrected_df
                st.session_state.correction_ledger = ledger_df
                st.session_state.corrected_process_map = corrected_process_map
                st.session_state.corrected_risk_report = corrected_risk_report

                corrected_docs_zip, corrected_docs_count = corrected_documents_zip_bytes(
                    report_df,
                    ledger_df,
                    st.session_state.get("raw_documents", {}),
                    ai_config,
                )
                st.session_state.corrected_docs_zip = corrected_docs_zip
                st.session_state.corrected_docs_count = corrected_docs_count

                st.success(f"Corrections applied. {len(ledger_df)} correction record(s) added to the ledger.")
                if corrected_docs_count:
                    st.success(f"Corrected document drafts are ready to download: {corrected_docs_count} document(s).")

    elif qdf is not None and qdf.empty:
        st.success("No high-impact clarification questions were generated for this scope.")



# ============================================================
# Streamlit UI
# ============================================================

st.set_page_config(page_title="AI Knowledge Inventory", page_icon="🧭", layout="wide")

st.markdown("""
<style>
.block-container { padding-top: 1.8rem; max-width: 1400px; }
.hero {
  padding: 1.25rem 1.35rem;
  border-radius: 22px;
  background: linear-gradient(135deg, rgba(37,99,235,0.12), rgba(14,165,233,0.08), rgba(255,255,255,0.9));
  border: 1px solid rgba(148, 163, 184, 0.28);
  margin-bottom: 1rem;
}
.hero-title { font-size: 2rem; font-weight: 800; letter-spacing: -0.03em; margin-bottom: 0.25rem; }
.hero-subtitle { color: #64748b; font-size: 1rem; }
.metric-card {
  padding: 1rem 1.1rem;
  border-radius: 18px;
  border: 1px solid rgba(148, 163, 184, 0.28);
  background: rgba(255, 255, 255, 0.82);
  box-shadow: 0 8px 24px rgba(15,23,42,0.05);
}
.metric-label { font-size: 0.85rem; color: #64748b; margin-bottom: 0.35rem; }
.metric-value { font-size: 1.9rem; font-weight: 800; }
.section-card {
  padding: 1rem 1.15rem;
  border-radius: 18px;
  border: 1px solid rgba(148, 163, 184, 0.28);
  background: rgba(255, 255, 255, 0.82);
  margin-bottom: 1rem;
}
.small-muted { color: #64748b; font-size: 0.9rem; }
.precorrect-cta {
  margin-top: 1.25rem;
  margin-bottom: 1.25rem;
  padding: 1.35rem 1.45rem;
  border-radius: 22px;
  border: 1px solid rgba(37, 99, 235, 0.22);
  background: linear-gradient(135deg, rgba(37,99,235,0.12), rgba(20,184,166,0.10), rgba(255,255,255,0.95));
  box-shadow: 0 14px 34px rgba(15,23,42,0.08);
}
.precorrect-title { font-size: 1.35rem; font-weight: 800; color: #0f172a; margin-bottom: 0.35rem; }
.precorrect-copy { color: #475569; font-size: 0.98rem; line-height: 1.55; }
.guided-card {
  padding: 1rem 1.15rem;
  border-radius: 18px;
  border: 1px solid rgba(148, 163, 184, 0.28);
  background: rgba(248, 250, 252, 0.88);
  margin-bottom: 1rem;
}
.recommended-pill {
  display: inline-block;
  padding: 0.28rem 0.55rem;
  border-radius: 999px;
  background: rgba(20,184,166,0.12);
  color: #0f766e;
  font-weight: 700;
  font-size: 0.82rem;
  margin-bottom: 0.45rem;
}
</style>
""", unsafe_allow_html=True)

st.markdown("""
<div class="hero">
  <div class="hero-title">🧭 AI Knowledge Inventory</div>
  <div class="hero-subtitle">Upload messy enterprise documents, diagnose knowledge readiness, run guided pre-correction, and download corrected draft documents.</div>
</div>
""", unsafe_allow_html=True)

default_config = get_default_ai_config()

with st.sidebar:
    st.header("Model")

    provider = st.selectbox(
        "Provider",
        ["No API", "OpenAI", "Gemini", "Claude", "Other"],
        index=0,
        help="Use No API for UI testing. Use Other for OpenAI-compatible providers."
    )

    api_key = ""
    base_url = ""
    model = ""
    mode = "heuristic"

    if provider == "No API":
        mode = "heuristic"
        st.info("Rule-based fallback. Useful for interface testing.")
    elif provider == "OpenAI":
        mode = "openai_compatible"
        api_key = st.text_input("API key", type="password", value=default_config.get("api_key", ""))
        model = st.text_input("Model", value="gpt-4o-mini")
        base_url = "https://api.openai.com/v1"
    elif provider == "Gemini":
        mode = "gemini"
        api_key = st.text_input("API key", type="password", value=default_config.get("api_key", ""))
        model = st.text_input("Model", value="gemini-2.5-flash")
        base_url = "https://generativelanguage.googleapis.com"
    elif provider == "Claude":
        mode = "anthropic"
        api_key = st.text_input("API key", type="password", value=default_config.get("api_key", ""))
        model = st.text_input("Model", value="claude-3-5-haiku-latest")
        base_url = "https://api.anthropic.com"
    elif provider == "Other":
        mode = "openai_compatible"
        api_key = st.text_input("API key", type="password", value=default_config.get("api_key", ""))
        base_url = st.text_input("Base URL", placeholder="https://api.deepseek.com/v1")
        model = st.text_input("Model", placeholder="deepseek-chat")
        st.caption("For DeepSeek, Qwen-compatible endpoints, Moonshot, local gateways, etc.")

    ai_config = {
        "mode": mode,
        "api_key": api_key,
        "base_url": base_url,
        "model": model,
    }

    st.markdown("---")
    st.header("Extraction")
    max_chars = st.slider("Text limit per document", 2000, 50000, 15000, step=1000)
    ocr_images = st.checkbox("Enable OCR for images/scanned PDFs", value=False)
    st.caption("OCR requires Tesseract. If unavailable, the app will continue and show extraction warnings.")

for key, default in {
    "inventory": None,
    "process_map": "",
    "risk_report": "",
    "raw_docs_count": 0,
    "extraction_warnings": [],
    "corrected_inventory": None,
    "corrected_process_map": "",
    "corrected_risk_report": "",
    "correction_questions": None,
    "correction_ledger": None,
    "correction_scope_docs": [],
    "correction_scope_type": "",
    "show_precorrection": False,
    "raw_documents": {},
    "corrected_docs_zip": None,
    "corrected_docs_count": 0,
}.items():
    if key not in st.session_state:
        st.session_state[key] = default

def run_analysis(extracted_docs):
    if not extracted_docs:
        st.error("No readable documents selected.")
        return

    st.session_state.raw_docs_count = len(extracted_docs)
    # Store extracted text in session so v12 can generate corrected draft documents after pre-correction.
    # The original uploaded files are not overwritten.
    st.session_state.raw_documents = {
        doc.get("document_name", f"document_{idx}"): {
            "document_name": doc.get("document_name", ""),
            "source_path": doc.get("source_path", ""),
            "file_type": doc.get("file_type", ""),
            "text": doc.get("text", "")[:30000],
        }
        for idx, doc in enumerate(extracted_docs, start=1)
    }
    warnings = []
    results = []
    total = len(extracted_docs)

    progress = st.progress(0)
    status = st.empty()

    for idx, doc in enumerate(extracted_docs, start=1):
        status.info(f"Analyzing {idx}/{total}: {doc['document_name']}")
        if doc.get("extraction_warnings"):
            warnings.extend([f"{doc['document_name']}: {w}" for w in doc["extraction_warnings"]])

        metadata = {
            "document_name": doc["document_name"],
            "source_path": doc["source_path"],
            "file_type": doc["file_type"],
            "sheet_names": doc.get("sheet_names", []),
            "text_length": len(doc.get("text", "")),
            "extraction_warnings": doc.get("extraction_warnings", []),
        }

        results.append(
            analyze_document(
                doc.get("text", ""),
                metadata,
                max_chars=max_chars,
                ai_config=ai_config
            )
        )
        progress.progress(idx / total)

    status.info("Generating process map and risk report...")

    df = pd.DataFrame(results)
    records = df.to_dict(orient="records")
    process_map = generate_process_map(records, ai_config=ai_config)
    risk_report = generate_risk_report(records, ai_config=ai_config)

    st.session_state.inventory = df
    st.session_state.process_map = process_map
    st.session_state.risk_report = risk_report
    st.session_state.extraction_warnings = warnings

    # Reset pre-correction state when a new analysis is run.
    st.session_state.corrected_inventory = None
    st.session_state.corrected_process_map = ""
    st.session_state.corrected_risk_report = ""
    st.session_state.correction_questions = None
    st.session_state.correction_ledger = None
    st.session_state.correction_scope_docs = []
    st.session_state.corrected_docs_zip = None
    st.session_state.corrected_docs_count = 0
    st.session_state.show_precorrection = False

    status.success("Initial analysis complete. Review the dashboard first, then start guided pre-correction from the large button below.")
    progress.progress(1.0)

def metric_card(label, value):
    st.markdown(f"""
    <div class="metric-card">
      <div class="metric-label">{label}</div>
      <div class="metric-value">{value}</div>
    </div>
    """, unsafe_allow_html=True)

def render_dashboard(df: pd.DataFrame):
    if df is None or df.empty:
        return

    total = len(df)
    high = int((df["risk_level"] == "High").sum()) if "risk_level" in df else 0
    medium = int((df["risk_level"] == "Medium").sum()) if "risk_level" in df else 0
    review = int(df["needs_human_review"].astype(bool).sum()) if "needs_human_review" in df else 0
    kb = int((df["can_enter_kb"] == "Yes").sum()) if "can_enter_kb" in df else 0

    cols = st.columns(5)
    with cols[0]: metric_card("Documents", total)
    with cols[1]: metric_card("High risk", high)
    with cols[2]: metric_card("Medium risk", medium)
    with cols[3]: metric_card("Need review", review)
    with cols[4]: metric_card("KB candidates", kb)

    st.markdown("")
    st.markdown('<div class="section-card">', unsafe_allow_html=True)
    st.subheader("Interactive dashboard")

    f1, f2, f3 = st.columns(3)
    risk_options = sorted(df["risk_level"].dropna().unique().tolist()) if "risk_level" in df else []
    dept_options = sorted(df["department"].dropna().unique().tolist()) if "department" in df else []
    kb_options = sorted(df["can_enter_kb"].dropna().unique().tolist()) if "can_enter_kb" in df else []

    risk_filter = f1.multiselect("Risk level", risk_options, default=risk_options)
    dept_filter = f2.multiselect("Department", dept_options, default=dept_options)
    kb_filter = f3.multiselect("KB readiness", kb_options, default=kb_options)

    filtered = df.copy()
    if risk_filter and "risk_level" in filtered:
        filtered = filtered[filtered["risk_level"].isin(risk_filter)]
    if dept_filter and "department" in filtered:
        filtered = filtered[filtered["department"].isin(dept_filter)]
    if kb_filter and "can_enter_kb" in filtered:
        filtered = filtered[filtered["can_enter_kb"].isin(kb_filter)]

    chart1, chart2 = st.columns(2)

    with chart1:
        if not filtered.empty and "risk_level" in filtered:
            order = ["High", "Medium", "Low"]
            counts = filtered["risk_level"].value_counts().reindex(order).dropna().reset_index()
            counts.columns = ["risk_level", "count"]
            fig = px.bar(
                counts,
                x="risk_level",
                y="count",
                title="Risk Distribution",
                text="count",
                color="risk_level",
                color_discrete_map={"High": "#ef4444", "Medium": "#f59e0b", "Low": "#22c55e"},
            )
            fig.update_layout(height=350, margin=dict(l=10, r=10, t=60, b=10))
            st.plotly_chart(fig, use_container_width=True)

    with chart2:
        if not filtered.empty and "department" in filtered:
            counts = filtered["department"].value_counts().head(10).reset_index()
            counts.columns = ["department", "count"]
            fig = px.bar(
                counts,
                x="count",
                y="department",
                orientation="h",
                title="Documents by Department",
                text="count",
                color="count",
                color_continuous_scale="Blues",
            )
            fig.update_layout(height=350, margin=dict(l=10, r=10, t=60, b=10), yaxis={"categoryorder": "total ascending"})
            st.plotly_chart(fig, use_container_width=True)

    chart3, chart4 = st.columns(2)

    with chart3:
        if not filtered.empty and "document_type" in filtered:
            counts = filtered["document_type"].value_counts().head(8).reset_index()
            counts.columns = ["document_type", "count"]
            fig = px.pie(counts, names="document_type", values="count", title="Document Type Mix", hole=0.55)
            fig.update_layout(height=380, margin=dict(l=10, r=10, t=60, b=10))
            st.plotly_chart(fig, use_container_width=True)

    with chart4:
        if not filtered.empty and "can_enter_kb" in filtered:
            counts = filtered["can_enter_kb"].value_counts().reset_index()
            counts.columns = ["can_enter_kb", "count"]
            fig = px.pie(
                counts,
                names="can_enter_kb",
                values="count",
                title="Knowledge Base Readiness",
                hole=0.55,
                color="can_enter_kb",
                color_discrete_map={"Yes": "#22c55e", "Needs Review": "#f59e0b", "No": "#ef4444"},
            )
            fig.update_layout(height=380, margin=dict(l=10, r=10, t=60, b=10))
            st.plotly_chart(fig, use_container_width=True)

    st.markdown("</div>", unsafe_allow_html=True)

    st.markdown('<div class="section-card">', unsafe_allow_html=True)
    st.subheader("Priority Review Items")
    review_df = filtered.copy()
    if "needs_human_review" in review_df:
        review_df = review_df[review_df["needs_human_review"].astype(bool)]
    if "risk_level" in review_df:
        review_df["_risk_order"] = review_df["risk_level"].map({"High": 0, "Medium": 1, "Low": 2}).fillna(3)
        review_df = review_df.sort_values("_risk_order").drop(columns=["_risk_order"])

    cols = ["document_name", "department", "business_process", "risk_level", "risk_reason", "recommended_action"]
    cols = [c for c in cols if c in review_df.columns]
    if not review_df.empty:
        st.dataframe(review_df[cols], use_container_width=True, hide_index=True)
    else:
        st.success("No priority review items in the current filter.")
    st.markdown("</div>", unsafe_allow_html=True)

    if st.session_state.extraction_warnings:
        with st.expander("Extraction warnings"):
            for warning in st.session_state.extraction_warnings:
                st.write(f"- {warning}")

st.markdown('<div class="section-card">', unsafe_allow_html=True)
st.subheader("1. Add documents")
st.markdown('<div class="small-muted">Recommended for folder upload: compress the folder into ZIP and upload it here.</div>', unsafe_allow_html=True)

tab_zip, tab_files = st.tabs(["Upload folder ZIP", "Upload files"])

zip_file = None
uploaded_files = None

with tab_zip:
    zip_file = st.file_uploader("Upload a ZIP file", type=["zip"], accept_multiple_files=False)
    if zip_file:
        try:
            preview_docs = extract_docs_from_zip_upload(zip_file, ocr_images=ocr_images)
            st.success(f"{len(preview_docs)} supported document(s) found in ZIP.")
        except Exception as exc:
            st.error(f"Could not read ZIP: {exc}")

with tab_files:
    uploaded_files = st.file_uploader(
        "Upload individual files",
        type=[ext.replace(".", "") for ext in SUPPORTED_EXTENSIONS],
        accept_multiple_files=True
    )
    if uploaded_files:
        st.success(f"{len(uploaded_files)} file(s) uploaded.")

st.markdown("</div>", unsafe_allow_html=True)

selected_count = 0
if zip_file:
    try:
        selected_count += len(extract_docs_from_zip_upload(zip_file, ocr_images=ocr_images))
    except Exception:
        pass
if uploaded_files:
    selected_count += len(uploaded_files)

st.markdown('<div class="section-card">', unsafe_allow_html=True)
st.subheader("2. Run analysis")
if selected_count:
    st.info(f"{selected_count} document(s) selected.")
else:
    st.warning("No documents selected yet.")

if st.button("Run Knowledge Inventory Analysis", type="primary", use_container_width=True):
    extracted = []
    if zip_file:
        extracted.extend(extract_docs_from_zip_upload(zip_file, ocr_images=ocr_images))
    if uploaded_files:
        extracted.extend([extract_text_from_upload(f, ocr_images=ocr_images) for f in uploaded_files])

    seen = set()
    unique_docs = []
    for doc in extracted:
        key = (doc.get("document_name"), doc.get("source_path"))
        if key not in seen:
            seen.add(key)
            unique_docs.append(doc)

    run_analysis(unique_docs)
st.markdown("</div>", unsafe_allow_html=True)


if st.session_state.inventory is not None:
    st.markdown("---")

    active_df = st.session_state.corrected_inventory if st.session_state.corrected_inventory is not None else st.session_state.inventory
    active_process_map = st.session_state.corrected_process_map if st.session_state.corrected_inventory is not None else st.session_state.process_map
    active_risk_report = st.session_state.corrected_risk_report if st.session_state.corrected_inventory is not None else st.session_state.risk_report

    if st.session_state.corrected_inventory is not None:
        st.success("Showing corrected dashboard after pre-correction.")
        dashboard_df = corrected_view_for_report(active_df)
    else:
        dashboard_df = active_df

    render_dashboard(dashboard_df)

    st.markdown("""
    <div class="precorrect-cta">
      <div class="precorrect-title">Ready to pre-correct and generate corrected documents?</div>
      <div class="precorrect-copy">
        After reviewing the dashboard, run a guided pre-correction round. The system will ask targeted questions, recommend the safest option, leave room for your own answer, and generate corrected draft documents plus a correction ledger after applying your answers.
      </div>
    </div>
    """, unsafe_allow_html=True)

    if st.button("Start Guided Pre-correction", type="primary", use_container_width=True):
        st.session_state.show_precorrection = True

    if st.session_state.show_precorrection:
        base_for_correction = st.session_state.corrected_inventory if st.session_state.corrected_inventory is not None else st.session_state.inventory
        base_for_correction = corrected_view_for_report(base_for_correction)
        render_precorrection_workbench(base_for_correction, ai_config)

        if st.session_state.correction_ledger is not None and not st.session_state.correction_ledger.empty:
            st.markdown("### Correction Ledger")
            st.dataframe(st.session_state.correction_ledger, use_container_width=True, hide_index=True)

            if st.session_state.corrected_docs_zip:
                st.success(f"Corrected documents are ready to download: {st.session_state.corrected_docs_count} document(s).")
                st.download_button(
                    "Download Corrected Documents",
                    data=st.session_state.corrected_docs_zip,
                    file_name="corrected_documents.zip",
                    mime="application/zip",
                    use_container_width=True,
                )

    st.markdown("---")
    tab1, tab2, tab3, tab4 = st.tabs([
        "Inventory table",
        "Process map",
        "Risk report",
        "Download reports"
    ])

    with tab1:
        display_df = corrected_view_for_report(active_df)
        display_cols = [
            "document_name", "file_type", "document_type", "department", "business_process",
            "knowledge_category", "risk_level", "can_enter_kb",
            "needs_human_review", "summary", "recommended_action"
        ]
        if st.session_state.corrected_inventory is not None:
            display_cols.extend(["risk_level_original", "can_enter_kb_original", "correction_notes"])
        display_cols = [c for c in display_cols if c in display_df.columns]
        st.dataframe(display_df[display_cols], use_container_width=True, hide_index=True)

    with tab2:
        st.markdown(active_process_map)

    with tab3:
        st.markdown(active_risk_report)

    with tab4:
        df_to_download = corrected_view_for_report(active_df)
        process_map = active_process_map
        risk_report = active_risk_report
        ledger_df = st.session_state.correction_ledger

        d1, d2, d3, d4, d5, d6 = st.columns(6)
        with d1:
            st.download_button(
                "Download Excel",
                data=inventory_excel_bytes(df_to_download),
                file_name="knowledge_inventory_corrected.xlsx" if st.session_state.corrected_inventory is not None else "knowledge_inventory.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True,
            )
        with d2:
            st.download_button(
                "Download Process Map",
                data=markdown_bytes(process_map),
                file_name="process_map_corrected.md" if st.session_state.corrected_inventory is not None else "process_map.md",
                mime="text/markdown",
                use_container_width=True,
            )
        with d3:
            st.download_button(
                "Download Risk Report",
                data=markdown_bytes(risk_report),
                file_name="document_risk_report_corrected.md" if st.session_state.corrected_inventory is not None else "document_risk_report.md",
                mime="text/markdown",
                use_container_width=True,
            )
        with d4:
            if ledger_df is not None and not ledger_df.empty:
                st.download_button(
                    "Download Ledger",
                    data=dataframe_excel_bytes(ledger_df, "Correction Ledger"),
                    file_name="correction_ledger.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True,
                )
            else:
                st.button("Download Ledger", disabled=True, use_container_width=True)

        with d5:
            if st.session_state.corrected_docs_zip:
                st.download_button(
                    "Download Corrected Docs",
                    data=st.session_state.corrected_docs_zip,
                    file_name="corrected_documents.zip",
                    mime="application/zip",
                    use_container_width=True,
                )
            else:
                st.button("Download Corrected Docs", disabled=True, use_container_width=True)

        with d6:
            st.download_button(
                "Download All",
                data=report_zip_bytes_with_ledger(df_to_download, process_map, risk_report, ledger_df, st.session_state.corrected_docs_zip),
                file_name="knowledge_inventory_corrected_bundle.zip" if st.session_state.corrected_inventory is not None else "knowledge_inventory_report_bundle.zip",
                mime="application/zip",
                use_container_width=True,
            )
